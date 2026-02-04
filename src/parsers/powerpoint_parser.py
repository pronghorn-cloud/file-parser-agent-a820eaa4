"""PowerPoint presentation parser using python-pptx."""

from pathlib import Path
from typing import Any, Dict, List
import logging
import base64
import io

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from .base import BaseDocumentParser, ParsedDocument

logger = logging.getLogger(__name__)


class PowerPointParser(BaseDocumentParser):
    """Parser for PowerPoint presentations.
    
    Extracts slides, text, titles, notes, shapes, images, and charts
    from PowerPoint files using the python-pptx library.
    """
    
    SUPPORTED_EXTENSIONS = ['.pptx', '.ppt']
    
    def __init__(self, file_path: Path):
        if not HAS_PPTX:
            raise ImportError("python-pptx is required for PowerPoint parsing. Install with: pip install python-pptx")
        super().__init__(file_path)
        self._presentation = None
        
        if self.file_path.suffix.lower() == '.ppt':
            logger.warning(".ppt files may have limited support. Consider converting to .pptx")
    
    @property
    def presentation(self) -> 'Presentation':
        """Lazy-load PowerPoint presentation."""
        if self._presentation is None:
            self._presentation = Presentation(str(self.file_path))
        return self._presentation
    
    def parse(self) -> ParsedDocument:
        """Parse PowerPoint and return structured content."""
        logger.info(f"Parsing PowerPoint: {self.file_path}")
        
        doc = ParsedDocument(
            filename=self.file_path.name,
            file_type='powerpoint'
        )
        
        try:
            doc.metadata = self.extract_metadata()
            doc.content = {
                'slide_count': len(self.presentation.slides),
                'slides': self._extract_slides()
            }
            doc.tables = self.extract_tables()
            doc.images = self.extract_images()
        except Exception as e:
            logger.error(f"Error parsing PowerPoint: {e}")
            doc.errors.append(str(e))
        
        return doc
    
    def extract_text(self) -> str:
        """Extract all text from PowerPoint."""
        text_parts = []
        for i, slide in enumerate(self.presentation.slides, start=1):
            text_parts.append(f"=== Slide {i} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)
            
            # Include speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    text_parts.append(f"[Speaker Notes: {notes_text}]")
            
            text_parts.append("")  # Empty line between slides
        
        return "\n".join(text_parts)
    
    def extract_metadata(self) -> Dict[str, Any]:
        """Extract PowerPoint metadata."""
        metadata = {}
        core_props = self.presentation.core_properties
        
        if core_props:
            prop_names = [
                'author', 'category', 'comments', 'content_status',
                'created', 'identifier', 'keywords', 'language',
                'last_modified_by', 'last_printed', 'modified',
                'revision', 'subject', 'title', 'version'
            ]
            
            for prop in prop_names:
                value = getattr(core_props, prop, None)
                if value:
                    if hasattr(value, 'isoformat'):
                        value = value.isoformat()
                    metadata[prop] = str(value)
        
        # Add presentation statistics
        metadata['slide_count'] = len(self.presentation.slides)
        
        # Get slide dimensions
        metadata['slide_width'] = self.presentation.slide_width.inches if self.presentation.slide_width else None
        metadata['slide_height'] = self.presentation.slide_height.inches if self.presentation.slide_height else None
        
        return metadata
    
    def _extract_slides(self) -> List[Dict[str, Any]]:
        """Extract content from all slides."""
        slides = []
        for i, slide in enumerate(self.presentation.slides, start=1):
            slide_data = {
                'slide_number': i,
                'layout': slide.slide_layout.name if slide.slide_layout else None,
                'title': self._get_slide_title(slide),
                'content': [],
                'shapes': [],
                'notes': None,
                'images': [],
                'tables': [],
                'charts': []
            }
            
            # Extract content from shapes
            for shape in slide.shapes:
                shape_info = self._extract_shape_info(shape)
                slide_data['shapes'].append(shape_info)
                
                # Categorize content
                if shape_info.get('text'):
                    slide_data['content'].append(shape_info['text'])
                if shape_info.get('is_image'):
                    slide_data['images'].append(shape_info)
                if shape_info.get('is_table'):
                    slide_data['tables'].append(shape_info.get('table_data'))
                if shape_info.get('is_chart'):
                    slide_data['charts'].append(shape_info)
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    slide_data['notes'] = notes_text
            
            slides.append(slide_data)
        
        return slides
    
    def _get_slide_title(self, slide) -> str:
        """Extract slide title."""
        if slide.shapes.title:
            return slide.shapes.title.text
        
        # Fallback: look for title placeholder
        for shape in slide.shapes:
            if shape.is_placeholder:
                if hasattr(shape, 'placeholder_format'):
                    if 'TITLE' in str(shape.placeholder_format.type):
                        return shape.text if hasattr(shape, 'text') else None
        
        return None
    
    def _extract_shape_info(self, shape) -> Dict[str, Any]:
        """Extract information from a shape."""
        info = {
            'name': shape.name,
            'shape_type': str(shape.shape_type) if hasattr(shape, 'shape_type') else None,
            'left': shape.left.inches if shape.left else None,
            'top': shape.top.inches if shape.top else None,
            'width': shape.width.inches if shape.width else None,
            'height': shape.height.inches if shape.height else None,
            'text': None,
            'is_image': False,
            'is_table': False,
            'is_chart': False
        }
        
        # Extract text
        if hasattr(shape, 'text') and shape.text.strip():
            info['text'] = shape.text
        
        # Check for image
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            info['is_image'] = True
            info['image_info'] = self._extract_image_info(shape)
        
        # Check for table
        if shape.has_table:
            info['is_table'] = True
            info['table_data'] = self._extract_table_from_shape(shape)
        
        # Check for chart
        if shape.has_chart:
            info['is_chart'] = True
            info['chart_info'] = self._extract_chart_info(shape)
        
        return info
    
    def _extract_image_info(self, shape) -> Dict[str, Any]:
        """Extract image information from a picture shape."""
        image_info = {
            'width': shape.width.inches if shape.width else None,
            'height': shape.height.inches if shape.height else None,
        }
        
        # Try to get image blob for AI analysis
        try:
            image = shape.image
            image_info['content_type'] = image.content_type
            image_info['ext'] = image.ext
            # Store base64 for potential vision API analysis
            image_info['blob'] = base64.b64encode(image.blob).decode('utf-8')
        except Exception as e:
            logger.debug(f"Could not extract image blob: {e}")
        
        return image_info
    
    def _extract_table_from_shape(self, shape) -> Dict[str, Any]:
        """Extract table data from a shape."""
        table = shape.table
        table_data = {
            'rows': len(table.rows),
            'columns': len(table.columns),
            'data': []
        }
        
        for row in table.rows:
            row_data = [cell.text for cell in row.cells]
            table_data['data'].append(row_data)
        
        return table_data
    
    def _extract_chart_info(self, shape) -> Dict[str, Any]:
        """Extract chart information from a shape."""
        chart = shape.chart
        chart_info = {
            'chart_type': str(chart.chart_type) if hasattr(chart, 'chart_type') else None,
            'has_title': chart.has_title,
            'title': chart.chart_title.text_frame.text if chart.has_title and chart.chart_title else None
        }
        
        return chart_info
    
    def extract_tables(self) -> List[Dict[str, Any]]:
        """Extract all tables from presentation."""
        tables = []
        for slide_num, slide in enumerate(self.presentation.slides, start=1):
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = self._extract_table_from_shape(shape)
                    table_data['slide_number'] = slide_num
                    table_data['shape_name'] = shape.name
                    tables.append(table_data)
        
        return tables
    
    def extract_images(self) -> List[Dict[str, Any]]:
        """Extract all image references from presentation."""
        images = []
        for slide_num, slide in enumerate(self.presentation.slides, start=1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_info = self._extract_image_info(shape)
                    image_info['slide_number'] = slide_num
                    image_info['shape_name'] = shape.name
                    images.append(image_info)
        
        return images
